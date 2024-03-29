from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import csv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
import spacy
import logging
import os
import sys
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN

# Create a new logger instance for decision-bubbles.py
logger = logging.getLogger('decision-bubbles')
logger.setLevel(logging.INFO)  # Set the log level to INFO

# Create a file handler to output logs to a file
file_handler = logging.FileHandler('decision-bubbles.log')
file_handler.setLevel(logging.INFO)  # Set the log level for the handler to INFO

# Define a formatter for the log messages
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)  # Attach the formatter to the handler

# Add the file handler to the logger
logger.addHandler(file_handler)

# Load the spaCy English model
nlp = spacy.load("en_core_web_sm")

try:
    # Get the file name with the identifier from the command-line arguments
    filename_with_identifier = sys.argv[1]
    logger.info("looking for the filepath")

    filename_without_extension = os.path.splitext(filename_with_identifier)[0]
    logger.info(f"Filename without extension: {filename_without_extension}")

    # Construct the complete file path
    file_path = os.path.join(filename_with_identifier)
    logger.info(f"Uploaded filepath: {file_path}")

    with open(file_path, 'r') as csv_file:
        csv_reader = csv.reader(csv_file)
        rows = list(csv_reader)

    # Create a PowerPoint presentation
    presentation = Presentation()

    # Create a blank slide
    slide_layout = presentation.slide_layouts[5]  # Blank slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Define initial positions for ovals
    left = Inches(0.5)
    top = Inches(5)

    # Iterate over all cells in the CSV file
    for row in rows:
        for idx, cell in enumerate(row):
            # Skip processing cells in columns A and F (indexes 0 and 5 respectively)
            if idx in [0, 5]:
                continue

            # Check if the cell begins with "or:" case-insensitive
            if cell.lower().startswith("or:"):
                # Extract the phrase after "or:"
                phrase = cell[3:].strip()

                # Check if the phrase contains parentheses
                if '(' in phrase and ')' in phrase:
                    # Split the phrase into main and inner phrases
                    main_phrase, inner_phrase = phrase.split('(', 1)
                    inner_phrase = inner_phrase[:-1]  # Remove the closing parenthesis

                    # Create the outer oval for the main phrase
                    outer_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(3), Inches(1.5))
                    outer_oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                    outer_oval.line.width = Pt(1)  # Border width
                    outer_oval.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # Dotted line border
                    outer_oval.fill.solid()
                    outer_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                    # Add the main phrase to the outer oval
                    text_frame_outer = outer_oval.text_frame
                    p_outer = text_frame_outer.add_paragraph()
                    p_outer.text = main_phrase.strip()
                    p_outer.font.size = Pt(14)  # Adjust font size as needed
                    p_outer.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                    p_outer.alignment = PP_ALIGN.CENTER  # Center align text

                    # Create the inner oval for the phrase inside parentheses
                    inner_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.5), top + Inches(0.3), Inches(2), Inches(0.8))
                    inner_oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                    inner_oval.line.width = Pt(1)  # Border width
                    inner_oval.fill.solid()
                    inner_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                    # Add the inner phrase to the inner oval
                    text_frame_inner = inner_oval.text_frame
                    p_inner = text_frame_inner.add_paragraph()
                    p_inner.text = inner_phrase.strip()
                    p_inner.font.size = Pt(12)  # Adjust font size as needed
                    p_inner.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                    p_inner.alignment = PP_ALIGN.CENTER  # Center align text

                else:
                    # Create a black dotted line oval for the extracted phrase
                    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(3), Inches(1.5))
                    oval.line.color.rgb = RGBColor(0, 0, 0)  # Black color for the border
                    oval.line.width = Pt(1)  # Border width
                    oval.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # Dotted line border
                    oval.fill.solid()
                    oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

                    # Add the extracted phrase to the oval
                    text_frame = oval.text_frame
                    p = text_frame.add_paragraph()
                    p.text = phrase.strip()
                    p.font.size = Pt(14)  # Adjust font size as needed
                    p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                    p.alignment = PP_ALIGN.CENTER  # Center align text

                # Update positions for the next oval
                left += Inches(3.5)
                if left + Inches(3) > Inches(10):
                    left = Inches(0.5)
                    top += Inches(2)

    # Save the PowerPoint presentation with the same identifier
    pptx_filename = os.path.join(filename_without_extension + '_decision-bubbles' + '.pptx')
    presentation.save(pptx_filename)
    logger.info("output_presentation is saved successfully.")
    logger.info(f"output_presentation path: {os.path.join(filename_without_extension +'_decision-bubbles'+'.pptx')}")
    print(f'CSV file has been converted to an editable PowerPoint presentation: "{pptx_filename}"')

except Exception as e:
    # Log any exceptions that occur
    logging.error(f'An error occurred: {str(e)}')
