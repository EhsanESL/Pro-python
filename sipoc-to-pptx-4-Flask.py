from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import csv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import spacy
import logging
import os
import sys

# Create a new logger instance for Seperate_verbs.py
logger = logging.getLogger('sipoc_to_pptx_logger')
logger.setLevel(logging.INFO)  # Set the log level to INFO

# Create a file handler to output logs to a file
file_handler = logging.FileHandler('sipoc_to_pptx.log')
file_handler.setLevel(logging.INFO)  # Set the log level for the handler to INFO

# Define a formatter for the log messages
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)  # Attach the formatter to the handler

# Add the file handler to the logger
logger.addHandler(file_handler)

# Load the spaCy English model
nlp = spacy.load("en_core_web_sm")

try:
    # Get the file name with the identifier from the command-line arguments
    filename_with_identifier = sys.argv[1]
    logger.info("looking for the filepath")

    filename_without_extension = os.path.splitext(filename_with_identifier)[0]
    logger.info(f"Filename without extension: {filename_without_extension}")

    # Construct the complete file path
    file_path = os.path.join(filename_with_identifier)
    logger.info(f"Uploaded filepath: {file_path}")

    with open(file_path, 'r') as csv_file:
        csv_reader = csv.reader(csv_file)
        rows = list(csv_reader)

    # Create a PowerPoint presentation
    presentation = Presentation()

    # Create a blank slide
    slide_layout = presentation.slide_layouts[5]  # Blank slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Define position and size for oval shapes
    num_cells = len(rows[1])
    oval_width = Inches(2)
    oval_height = Inches(0.8)
    text_font_size = Pt(18)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Calculate the diagonal line positions
    diagonal_length = (slide_width ** 2 + slide_height ** 2) ** 0.5
    x_step = slide_width / num_cells
    y_step = slide_height / num_cells

   
    # Create oval shapes and add text to them with specified font and background colors
    for i, cell in enumerate(rows[1]):
        # Skip column D
        if i == 3:
            continue
        
        left = i * x_step
        top = i * y_step
        if i == 0 or i == 5:  # First cell in column A and F
            # Create rectangle shape instead of oval
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for rectangles
        elif i == 1 or i == 4:  # Column B and E
            # Create oval shape with green background
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background
        elif i == 2:  # Column C
            # Create oval shape with red background
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, oval_width, oval_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background

        text_frame = shape.text_frame
        text_frame.text = cell
        text_frame.paragraphs[0].font.size = text_font_size
        text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


        # Use spaCy to extract verbs from the second row
        verbs = []
        for cell in rows[1]:
            doc = nlp(cell)
            verbs.extend([token.text for token in doc if token.pos_ == "VERB"])

        # Create nodes for verbs with no border and background at the bottom of the slide
        for i, verb in enumerate(verbs):
            left = i * x_step
            top = slide_height - oval_height - Inches(0.2)
            verb_oval = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, oval_width, oval_height)
            verb_oval.line.fill.solid()
            verb_oval.line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # No border
            verb_oval.line.width = Pt(0)  # No border width
            verb_oval.shadow.inherit = False  # No shadow
            verb_oval.fill.solid()
            verb_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)  # No background
            verb_text_frame = verb_oval.text_frame
            p = verb_text_frame.add_paragraph()
            p.text = verb
            p.font.size = text_font_size
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
            p.alignment = PP_ALIGN.CENTER

    # Save the PowerPoint presentation with the same identifier
    pptx_filename = os.path.join(filename_without_extension +'.pptx')
    presentation.save(pptx_filename)
    logger.info("output_presentation is saved successfully.")
    logger.info(f"output_presentation path: {os.path.join(filename_without_extension +'.pptx')}")
    print(f'Second row of CSV file has been converted to an editable PowerPoint presentation: "{pptx_filename}"')

except Exception as e:
    # Log any exceptions that occur
    logging.error(f'An error occurred: {str(e)}')

